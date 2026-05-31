"""Extract text, tables, and glossary matches from PPTX files."""

import os
import re

import openpyxl
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from config import FLAG_PATTERNS, FLAG_SHORT_TEXT, GLOSSARY_FILE

KOREAN_RE = re.compile(r"[\uAC00-\uD7A3\u1100-\u11FF\u3130-\u318F]+")
LATIN_RE = re.compile(r"[a-zA-ZÀ-ÿ]")


def load_glossary():
    """Load the FR-KO glossary workbook as a dictionary."""
    glossary = {}
    if not os.path.exists(GLOSSARY_FILE):
        print(f"[warning] Glossary file not found: {GLOSSARY_FILE}")
        return glossary

    wb = openpyxl.load_workbook(GLOSSARY_FILE, data_only=True)
    ws = wb.active
    for fr, ko, *_ in ws.iter_rows(min_row=2, values_only=True):
        if fr and ko:
            glossary[str(fr).strip()] = str(ko).strip()
    print(f"[glossary] Loaded {len(glossary)} terms")
    return glossary


def apply_glossary(text, glossary):
    """Return glossary matches found in the source text."""
    matched = []
    lowered = text.lower()
    for fr in sorted(glossary, key=len, reverse=True):
        if fr.lower() in lowered:
            matched.append(f"{fr} -> {glossary[fr]}")
    return matched


def flag_text(text):
    """Return simple review flags for short or suspicious text."""
    flags = []
    words = text.strip().split()
    if 0 < len(words) <= FLAG_SHORT_TEXT:
        flags.append(f"short text ({len(words)} words)")
    for pattern in FLAG_PATTERNS:
        if re.search(pattern, text):
            flags.append(f"pattern:{pattern}")
    return " | ".join(flags)


def strip_korean_annotations(text):
    """Remove Korean characters embedded in French source text."""
    if not LATIN_RE.search(text):
        return text
    cleaned = KOREAN_RE.sub("", text)
    cleaned = re.sub(r"  +", " ", cleaned).strip()
    return cleaned


def safe_shape_name(shape, index):
    """Return the PowerPoint shape name, or a stable fallback."""
    return shape.name or f"shape_{index}"


def iter_shapes_with_paths(shapes, parent_path=""):
    """Yield shapes recursively with parent/child group-path names."""
    for index, shape in enumerate(shapes):
        name = safe_shape_name(shape, index)
        shape_name = f"{parent_path}/{name}" if parent_path else name
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_with_paths(shape.shapes, shape_name)
        else:
            yield shape, shape_name


def get_run_texts(text_frame):
    """Return non-empty run texts from a text frame."""
    runs = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                runs.append(run.text)
    return runs


def make_item(slide_num, shape, shape_name, shape_type, text, glossary, run_texts=None):
    """Build one extracted text record."""
    full_text = strip_korean_annotations(text).strip()
    return {
        "slide_num": slide_num,
        "shape_id": getattr(shape, "shape_id", None),
        "shape_name": shape_name,
        "shape_type": shape_type,
        "run_texts": run_texts if run_texts is not None else [full_text],
        "full_text": full_text,
        "matched_terms": apply_glossary(full_text, glossary),
        "flag": flag_text(full_text),
    }


def extract_shape_items(slide_num, shape, shape_name, glossary):
    """Extract textbox and table-cell records from one shape."""
    items = []

    if shape.has_text_frame:
        full_text = shape.text_frame.text.strip()
        if full_text:
            items.append(
                make_item(
                    slide_num,
                    shape,
                    shape_name,
                    "textbox",
                    full_text,
                    glossary,
                    get_run_texts(shape.text_frame),
                )
            )

    if shape.has_table:
        for row_idx, row in enumerate(shape.table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_text = cell.text.strip()
                if not cell_text:
                    continue
                cell_name = f"{shape_name}_row{row_idx}_col{col_idx}"
                items.append(
                    make_item(
                        slide_num,
                        shape,
                        cell_name,
                        "table_cell",
                        cell_text,
                        glossary,
                    )
                )

    return items


def extract_from_pptx(pptx_path):
    """Extract text records from all slides in a PPTX file."""
    glossary = load_glossary()
    prs = Presentation(pptx_path)
    results = []

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape, shape_name in iter_shapes_with_paths(slide.shapes):
            results.extend(extract_shape_items(slide_num, shape, shape_name, glossary))

    print(f"[extract] {len(results)} text items from {len(prs.slides)} slides")
    return results, prs
