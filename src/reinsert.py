"""Reinsert reviewed Korean translations into the original PPTX."""

import os
import sys

import openpyxl
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from config import FINAL_DIR, INPUT_DIR, REVIEWED_DIR


def safe_shape_name(shape, index):
    """Return the PowerPoint shape name, or a stable fallback."""
    return shape.name or f"shape_{index}"


def iter_shapes_with_paths(shapes, parent_path=""):
    """Yield shapes recursively with the same group paths as extractor.py."""
    for index, shape in enumerate(shapes):
        name = safe_shape_name(shape, index)
        shape_name = f"{parent_path}/{name}" if parent_path else name
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_with_paths(shape.shapes, shape_name)
        else:
            yield shape, shape_name


def header_pos(header, name):
    """Return a required Excel header position."""
    if name not in header:
        raise ValueError(f"Missing required column: {name}")
    return header.index(name)


def load_translations(xlsx_path):
    """Load final translations keyed by slide, shape name, and shape type."""
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    ws = wb.active
    header = [cell.value for cell in ws[1]]

    i_slide = header_pos(header, "슬라이드")
    i_shape_name = header_pos(header, "Shape 이름")
    i_type = header_pos(header, "유형")
    i_final = header_pos(header, "최종 번역")

    translations = {}
    for row in ws.iter_rows(min_row=2, values_only=True):
        final = row[i_final]
        if final is None or not str(final).strip():
            continue
        key = (row[i_slide], row[i_shape_name], row[i_type])
        translations[key] = str(final)
    return translations


def set_textframe_text(text_frame, new_text):
    """Replace a text frame's text while preserving the first run style."""
    paragraphs = text_frame.paragraphs
    if not paragraphs:
        return

    first_para = paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = new_text
        for run in first_para.runs[1:]:
            run.text = ""
    else:
        first_para.text = new_text

    for para in paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def reinsert(xlsx_path, pptx_path):
    """Apply reviewed translations to the matching PPTX shapes."""
    translations = load_translations(xlsx_path)
    applied_keys = set()
    prs = Presentation(pptx_path)

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape, shape_name in iter_shapes_with_paths(slide.shapes):
            key = (slide_num, shape_name, "textbox")
            if shape.has_text_frame and key in translations:
                set_textframe_text(shape.text_frame, translations[key])
                applied_keys.add(key)

            if shape.has_table:
                for row_idx, table_row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(table_row.cells):
                        cell_name = f"{shape_name}_row{row_idx}_col{col_idx}"
                        key = (slide_num, cell_name, "table_cell")
                        if key in translations:
                            set_textframe_text(cell.text_frame, translations[key])
                            applied_keys.add(key)

    for slide_num, shape_name, _shape_type in sorted(set(translations) - applied_keys):
        print(f"[경고] 매칭 실패: slide {slide_num} / {shape_name}")

    base = os.path.splitext(os.path.basename(pptx_path))[0]
    out_path = os.path.join(FINAL_DIR, f"{base}_final.pptx")
    prs.save(out_path)
    print(f"[reinsert] Applied {len(applied_keys)} translations")
    print(f"[done] Saved: {out_path}")
    return out_path


def select_file(folder, extension, prompt):
    """Prompt the user to select a file from a folder."""
    files = sorted(f for f in os.listdir(folder) if f.lower().endswith(extension))
    if not files:
        print(f"[error] No {extension} files found in {folder}")
        return None

    print(f"\n{prompt}")
    for idx, filename in enumerate(files, 1):
        print(f"  {idx}. {filename}")

    try:
        return os.path.join(folder, files[int(input("Select number: ").strip()) - 1])
    except (ValueError, IndexError):
        print("[error] Invalid selection.")
        return None


def main():
    """Run the interactive reinsertion stage."""
    xlsx_path = select_file(REVIEWED_DIR, ".xlsx", "Select reviewed workbook from 3_reviewed:")
    if not xlsx_path:
        return
    pptx_path = select_file(INPUT_DIR, ".pptx", "Select original PPTX from input:")
    if not pptx_path:
        return
    reinsert(xlsx_path, pptx_path)


if __name__ == "__main__":
    main()
